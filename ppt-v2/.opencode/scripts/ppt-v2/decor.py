from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree
import theme


def add_orange_accent(slide):
    sp = slide.shapes.add_shape(
        9, Inches(12.7), Inches(0.3),
        Inches(0.4), Inches(0.4)
    )
    sp.fill.solid()
    sp.fill.fore_color.rgb = theme.BANNER_ORANGE
    sp.line.fill.background()
    sp.fill.fore_color.brightness = 0.85


def add_dotted_corner(slide):
    left = Inches(11.5)
    top = Inches(0.5)
    w = Inches(1.2)
    h = Inches(0.6)

    sp = slide.shapes.add_shape(
        1, left, top, w, h
    )
    sp.fill.background()
    sp.line.color.rgb = theme.BANNER_ORANGE
    sp.line.width = Pt(0.5)
    sp.line.dash_style = 2

    sp2 = slide.shapes.add_shape(
        1, left, int(top + h + Inches(0.08)),
        Inches(0.6), Pt(1)
    )
    sp2.fill.solid()
    sp2.fill.fore_color.rgb = theme.BORDER_RGB
    sp2.line.fill.background()


def add_geometric_dot(slide):
    sp = slide.shapes.add_shape(
        9, Inches(12.4), Inches(6.8),
        Inches(0.3), Inches(0.3)
    )
    sp.fill.solid()
    sp.fill.fore_color.rgb = theme.LIGHT_BG_RGB
    sp.line.fill.background()

    sp2 = slide.shapes.add_shape(
        9, Inches(12.5), Inches(6.9),
        Inches(0.15), Inches(0.15)
    )
    sp2.fill.solid()
    sp2.fill.fore_color.rgb = theme.BORDER_RGB
    sp2.line.fill.background()


def add_top_accent_line(slide):
    line = slide.shapes.add_shape(
        1, theme.MARGIN_L, Inches(0.25),
        Inches(2.5), Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = theme.BANNER_ORANGE
    line.line.fill.background()


def add_subtle_curve(slide):
    sp = slide.shapes.add_shape(
        1, Inches(11.0), Inches(0),
        Inches(2.333), Inches(1.2)
    )
    sp.fill.background()
    sp.line.color.rgb = theme.BORDER_RGB
    sp.line.width = Pt(0.5)
    sp.line.dash_style = 2

    sp2 = slide.shapes.add_shape(
        9, Inches(11.5), Inches(0.6),
        Inches(0.08), Inches(0.08)
    )
    sp2.fill.solid()
    sp2.fill.fore_color.rgb = theme.BANNER_ORANGE
    sp2.line.fill.background()


def add_design_elements(slide, count=1):
    if count <= 0:
        return
    add_orange_accent(slide)
    if count >= 2:
        add_top_accent_line(slide)
    if count >= 3:
        add_geometric_dot(slide)

    return True
