from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
import theme
import icons


def add_card(slide, left, top, width, height, icon_name, title, body_text,
             icon_size=None):
    if icon_size is None:
        icon_size = theme.SMALL_ICON_SIZE

    bg = slide.shapes.add_shape(
        5, left, top, width, height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = theme.LIGHT_BG_RGB
    bg.line.color.rgb = theme.BORDER_RGB
    bg.line.width = Pt(0.5)

    bg.shadow.inherit = False
    sp = bg._element
    spPr = sp.find(qn('p:spPr'))
    if spPr is not None:
        xmlns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        effLst = spPr.find(qn('a:effectLst'))
        if effLst is None:
            from lxml import etree
            effLst = etree.SubElement(spPr, qn('a:effectLst'))
        outerShdw = etree.SubElement(effLst, qn('a:outerShdw'))
        outerShdw.set('blurRad', str(Emu(60000)))
        outerShdw.set('dist', str(Emu(25400)))
        outerShdw.set('dir', str(18000000))
        outerShdw.set('algn', 'tl')
        srgbClr = etree.SubElement(outerShdw, qn('a:srgbClr'))
        srgbClr.set('val', '000000')
        alpha = etree.SubElement(srgbClr, qn('a:alpha'))
        alpha.set('val', '15000')

    icon_path = icons.get_icon_path(icon_name) if icon_name else None
    icon_top = top + height // 2 - icon_size // 2
    icon_left = left + Inches(0.25)

    if icon_path:
        try:
            slide.shapes.add_picture(icon_path, icon_left, icon_top,
                                     icon_size, icon_size)
        except Exception:
            pass

    content_left = left + icon_size + Inches(0.35)
    content_w = width - icon_size - Inches(0.55)

    title_box = slide.shapes.add_textbox(
        content_left, top + Inches(0.2), content_w, Inches(0.45)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    tf.margin_top = Pt(0)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = theme.DARK_RGB
    p.font.name = theme.FONT_FAMILY
    p.space_after = Pt(4)

    body_top = top + Inches(0.7)
    body_box = slide.shapes.add_textbox(
        content_left, body_top, content_w,
        height - Inches(0.85)
    )
    tf2 = body_box.text_frame
    tf2.word_wrap = True
    tf2.margin_left = Pt(0)
    tf2.margin_top = Pt(0)
    p2 = tf2.paragraphs[0]
    p2.text = body_text
    p2.font.size = Pt(14)
    p2.font.color.rgb = theme.BODY_COLOR
    p2.font.name = theme.FONT_FAMILY
    p2.space_after = Pt(0)
    p2.space_before = Pt(0)
    tf2.margin_left = Pt(0)
    tf2.margin_right = Pt(0)
    tf2.margin_bottom = Pt(0)
    tf2.word_wrap = True

    return bg
