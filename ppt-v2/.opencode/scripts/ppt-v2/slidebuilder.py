from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import theme
import layouts
import footer
import icons
import os


LAYOUT_MAP = {
    "cover": "cover_orange",
    "cover_orange": "cover_orange",
    "message": "message_only",
    "message_only": "message_only",
    "three_cards": "three_cards",
    "four_cards": "four_cards",
    "problem_solution": "problem_solution",
    "process": "process_horizontal",
    "process_horizontal": "process_horizontal",
    "kpi": "kpi_context",
    "kpi_context": "kpi_context",
    "lessons": "lessons_learned",
    "lessons_learned": "lessons_learned",
    "role": "role_focus",
    "role_focus": "role_focus",
    "adoption_loop": "adoption_loop",
    "loop": "adoption_loop",
    "closing": "closing",
    "conclusion": "closing",
}

SLIDE_TYPE_MAP = {
    "contexte": "kpi_context",
    "chiffres": "kpi_context",
    "problematique": "problem_solution",
    "probleme": "problem_solution",
    "question": "message_only",
    "reponse": "message_only",
    "processus": "process_horizontal",
    "roles": "role_focus",
    "enseignements": "lessons_learned",
    "bilan": "closing",
    "conclusion": "closing",
    "message": "message_only",
    "cartes": "three_cards",
    "adoption": "adoption_loop",
    "boucle": "adoption_loop",
}

CONCLUSION_TOPICS = {
    "contexte": "la comprehension du contexte est le point de depart de toute transformation",
    "problematique": "la resolution du probleme passe par une approche structuree",
    "question": "la reponse apportee eclaire la prise de decision",
    "processus": "la standardisation du processus garantit la reproductibilite",
    "roles": "la mobilisation des equipes est la cle du succes",
    "enseignements": "le partage des enseignements accelere la maturation",
    "chiffres": "la mesure de la performance pilote l'amelioration continue",
    "bilan": "l'amelioration continue est le moteur de la transformation",
    "adoption": "l'adoption reelle passe par l'appropriation par les utilisateurs",
    "default": "la valeur se cree dans la duree par l'engagement collectif",
}

DEFAULT_CONCLUSIONS = [
    "La cle du succes repose sur l'engagement de toutes les parties prenantes",
    "La transformation est un voyage qui se construit etape par etape",
    "L'innovation prend tout son sens quand elle est au service du metier",
    "La gouvernance est le pilier d'une transformation durable",
    "La mesure et le pilotage sont les boussoles de la reussite",
    "L'adoption reelle passe par l'appropriation par les utilisateurs",
    "La valeur se cree dans la duree par l'engagement collectif",
]


def _resolve_layout(slide_data):
    if slide_data.get("layout") and slide_data["layout"] in LAYOUT_MAP:
        return LAYOUT_MAP[slide_data["layout"]]
    slide_type = slide_data.get("type", "").lower()
    if slide_type in SLIDE_TYPE_MAP:
        return SLIDE_TYPE_MAP[slide_type]
    return "three_cards"


def _default_conclusion(slide_data):
    slide_type = slide_data.get("type", "").lower()
    topic = CONCLUSION_TOPICS.get(slide_type, CONCLUSION_TOPICS["default"])
    import random
    template = random.choice(DEFAULT_CONCLUSIONS)
    return template


def _normalize_cards(slide_data):
    if "cards" in slide_data and slide_data["cards"]:
        return slide_data["cards"]
    kpis = slide_data.get("kpis", [])
    if kpis:
        return [{"title": k.get("value", ""), "text": k.get("label", ""),
                 "icon": k.get("icon", "metrics")} for k in kpis]
    lessons = slide_data.get("lessons", [])
    if lessons:
        return [{"title": "", "text": l if isinstance(l, str) else l.get("text", ""),
                 "icon": "lightbulb"} for l in lessons]
    steps = slide_data.get("steps", [])
    if steps:
        return [{"title": s.get("title", ""), "text": s.get("description", s.get("text", "")),
                 "icon": "process"} for s in steps]
    phases = slide_data.get("phases", [])
    if phases:
        return [{"title": p.get("title", ""), "text": p.get("description", p.get("text", "")),
                 "icon": "adoption"} for p in phases]
    return []


def _normalize_kpis(slide_data):
    cards = _normalize_cards(slide_data)
    if cards:
        return [{"value": c.get("title", ""), "label": c.get("text", ""),
                 "icon": c.get("icon", "metrics")} for c in cards]
    return []


def _normalize_steps(slide_data):
    cards = _normalize_cards(slide_data)
    if cards:
        return [{"title": c.get("title", ""), "text": c.get("text", ""),
                 "icon": c.get("icon", "process")} for c in cards]
    return []


def build_deck(deck, output_path, icons_dir=None):
    icons.set_icons_dir(icons_dir) if icons_dir else None

    prs = Presentation()
    prs.slide_width = theme.SLIDE_W
    prs.slide_height = theme.SLIDE_H

    cover_data = deck.get("cover", {})
    cover_slide = prs.slides.add_slide(prs.slide_layouts[6])
    layouts.build_cover(cover_slide, cover_data)

    slides_data = deck.get("slides", [])
    total = len(slides_data) + 1

    for idx, slide_data in enumerate(slides_data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        layout_name = _resolve_layout(slide_data)

        if not slide_data.get("conclusion"):
            slide_data["conclusion"] = _default_conclusion(slide_data)
        if not slide_data.get("subtitle"):
            slide_data["subtitle"] = ""

        if layout_name == "kpi_context":
            slide_data["kpis"] = _normalize_kpis(slide_data)
        elif layout_name in ("process_horizontal",):
            slide_data["steps"] = _normalize_steps(slide_data)

        builder = layouts.BUILDERS.get(layout_name, layouts.build_three_cards)
        builder(slide, slide_data)

        footer.add_footer(slide, page_num=idx + 2, total=total)

    prs.save(output_path)
    return {"slides": total, "path": output_path}
