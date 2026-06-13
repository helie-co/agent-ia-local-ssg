from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import theme


def add_footer(slide, page_num=None, total=None):
    left = theme.MARGIN_L
    top = theme.FOOTER_Y
    width = theme.CONTENT_W
    height = theme.FOOTER_H

    sep_top = int(top - Inches(0.12))
    sep = slide.shapes.add_shape(
        1, left, sep_top, width, Pt(1)
    )
    sep.fill.solid()
    sep.fill.fore_color.rgb = theme.BORDER_RGB
    sep.line.fill.background()

    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)

    p = tf.paragraphs[0]
    p.font.size = theme.FOOTER_SIZE
    p.font.color.rgb = theme.MEDIUM_RGB
    p.font.name = theme.FONT_FAMILY
    p.alignment = PP_ALIGN.LEFT
    p.text = "C2 – Usage restreint"

    if page_num is not None:
        line = tf.add_paragraph()
        line.font.size = theme.FOOTER_SIZE
        line.font.color.rgb = theme.MEDIUM_RGB
        line.font.name = theme.FONT_FAMILY
        line.alignment = PP_ALIGN.RIGHT
        if total is not None:
            line.text = f"{page_num} / {total}"
        else:
            line.text = str(page_num)
