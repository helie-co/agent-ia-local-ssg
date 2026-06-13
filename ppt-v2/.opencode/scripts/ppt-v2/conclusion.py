from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import theme


def add_conclusion_banner(slide, text):
    left = Inches(0)
    top = theme.CONCLUSION_Y
    width = theme.SLIDE_W
    height = theme.CONCLUSION_H

    bg = slide.shapes.add_shape(1, left, top, width, height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = theme.BANNER_ORANGE_LIGHT
    bg.line.fill.background()

    bar = slide.shapes.add_shape(
        1, left, top, Inches(0.08), height
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = theme.BANNER_ORANGE
    bar.line.fill.background()

    box = slide.shapes.add_textbox(
        theme.MARGIN_L, top + Inches(0.12),
        theme.CONTENT_W, height - Inches(0.2)
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT

    run = p.add_run()
    run.text = "▎ "
    run.font.size = theme.CONCLUSION_SIZE
    run.font.color.rgb = theme.BANNER_ORANGE
    run.font.name = theme.FONT_FAMILY

    run2 = p.add_run()
    run2.text = text
    run2.font.size = theme.CONCLUSION_SIZE
    run2.font.bold = True
    run2.font.color.rgb = theme.DARK_RGB
    run2.font.name = theme.FONT_FAMILY

    p.space_before = Pt(0)
    p.space_after = Pt(0)
    for run in p.runs:
        run.font.language = None
