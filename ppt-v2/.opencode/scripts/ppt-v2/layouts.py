from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree
import theme
import card
import footer
import conclusion as conc
import decor


def _add_bg(slide, color_rgb):
    bg = slide.shapes.add_shape(
        1, 0, 0, theme.SLIDE_W, theme.SLIDE_H
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = color_rgb
    bg.line.fill.background()
    bg.shadow.inherit = False


def _add_title(slide, text, top=None, size=None):
    if top is None:
        top = theme.TITLE_TOP
    if size is None:
        size = theme.TITLE_SIZE
    box = slide.shapes.add_textbox(
        theme.MARGIN_L, top, theme.CONTENT_W, theme.TITLE_H
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    tf.margin_top = Pt(0)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.bold = True
    p.font.color.rgb = theme.DARK_RGB
    p.font.name = theme.FONT_FAMILY
    p.space_after = Pt(0)
    return box


def _add_subtitle(slide, text, top=None):
    if top is None:
        top = theme.SUBTITLE_TOP
    if not text:
        return None
    box = slide.shapes.add_textbox(
        theme.MARGIN_L, top, theme.CONTENT_W, theme.SUBTITLE_H
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    tf.margin_top = Pt(0)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = theme.SUBTITLE_SIZE
    p.font.color.rgb = theme.MEDIUM_RGB
    p.font.name = theme.FONT_FAMILY
    return box


def _add_kpi_big(slide, value, label, left, top, width):
    box = slide.shapes.add_textbox(left, top, width, Inches(1.2))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    tf.margin_top = Pt(0)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = value
    run.font.size = Pt(48)
    run.font.bold = True
    run.font.color.rgb = theme.BANNER_ORANGE
    run.font.name = theme.FONT_FAMILY

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(4)
    p2.space_after = Pt(0)
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = Pt(15)
    run2.font.color.rgb = theme.MEDIUM_RGB
    run2.font.name = theme.FONT_FAMILY


def _add_card_grid(slide, cards, cols=3):
    n = len(cards)
    if n == 0:
        return
    cols = min(cols, n)
    total_gap = theme.CARD_GAP * (cols - 1)
    card_w = (theme.CONTENT_W - total_gap) // cols
    rows = (n + cols - 1) // cols
    content_top = theme.CARD_Y
    card_h_adj = min(theme.CARD_H, (Inches(5.0) - (rows - 1) * Inches(0.4)) / rows)

    for i, c in enumerate(cards):
        col = i % cols
        row = i // cols
        left = theme.MARGIN_L + col * (card_w + theme.CARD_GAP)
        top = content_top + row * (card_h_adj + Inches(0.4))
        card.add_card(
            slide, left, top, card_w, card_h_adj,
            c.get("icon"), c.get("title", ""), c.get("text", ""),
        )


def _add_step_line_premium(slide, steps, top=None):
    if top is None:
        top = Inches(2.4)
    n = len(steps)
    if n == 0:
        return

    step_w = theme.CONTENT_W // n
    line_y = top + Inches(0.35)
    line = slide.shapes.add_shape(
        1, theme.MARGIN_L + Inches(0.4), line_y,
        theme.CONTENT_W - Inches(0.8), Pt(2)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = theme.BORDER_RGB
    line.line.fill.background()

    for i, s in enumerate(steps):
        cx = theme.MARGIN_L + step_w * i + step_w // 2

        circle = slide.shapes.add_shape(
            9, cx - Inches(0.2), line_y - Inches(0.2),
            Inches(0.4), Inches(0.4)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = theme.BANNER_ORANGE if i == 0 else theme.WHITE_RGB
        circle.line.color.rgb = theme.BANNER_ORANGE if i == 0 else theme.BORDER_RGB
        circle.line.width = Pt(1.5)

        nbox = slide.shapes.add_textbox(
            cx - Inches(0.2), line_y - Inches(0.2),
            Inches(0.4), Inches(0.4)
        )
        ntf = nbox.text_frame
        ntf.margin_left = Pt(0)
        ntf.margin_top = Pt(1)
        np = ntf.paragraphs[0]
        np.alignment = PP_ALIGN.CENTER
        np.text = str(i + 1)
        np.font.size = Pt(14)
        np.font.bold = True
        np.font.color.rgb = theme.WHITE_RGB if i == 0 else theme.DARK_RGB
        np.font.name = theme.FONT_FAMILY

        icon_name = s.get("icon", "process")
        icon_path = None
        icon_dir_holder = []
        import icons as ico_mod
        icon_path = ico_mod.get_icon_path(icon_name)
        if icon_path:
            try:
                slide.shapes.add_picture(
                    icon_path,
                    cx - Inches(0.25), line_y + Inches(0.35),
                    Inches(0.5), Inches(0.5)
                )
            except Exception:
                pass

        tbox = slide.shapes.add_textbox(
            theme.MARGIN_L + step_w * i + Inches(0.15),
            line_y + Inches(0.95),
            step_w - Inches(0.3), Inches(0.6)
        )
        ttf = tbox.text_frame
        ttf.word_wrap = True
        ttf.margin_left = Pt(0)
        ttf.margin_top = Pt(0)
        tp = ttf.paragraphs[0]
        tp.text = s.get("title", "")
        tp.font.size = Pt(16)
        tp.font.bold = True
        tp.font.color.rgb = theme.DARK_RGB
        tp.font.name = theme.FONT_FAMILY
        tp.alignment = PP_ALIGN.CENTER

        if s.get("text"):
            dbox = slide.shapes.add_textbox(
                theme.MARGIN_L + step_w * i + Inches(0.15),
                line_y + Inches(1.55),
                step_w - Inches(0.3), Inches(0.8)
            )
            dtf = dbox.text_frame
            dtf.word_wrap = True
            dtf.margin_left = Pt(0)
            dtf.margin_top = Pt(0)
            dp = dtf.paragraphs[0]
            dp.text = s["text"]
            dp.font.size = Pt(14)
            dp.font.color.rgb = theme.MEDIUM_RGB
            dp.font.name = theme.FONT_FAMILY
            dp.alignment = PP_ALIGN.CENTER
            dp.space_before = Pt(0)


def _add_cycle_premium(slide, phases, center_x=None, center_y=None, radius=None):
    if center_x is None:
        center_x = Inches(6.666)
    if center_y is None:
        center_y = Inches(3.5)
    if radius is None:
        radius = Inches(2.0)

    n = len(phases)
    if n == 0:
        return
    cx = int(center_x)
    cy = int(center_y)

    import math
    for i, p in enumerate(phases):
        angle = 2 * math.pi * i / n - math.pi / 2
        px = int(cx + radius * math.cos(angle))
        py = int(cy + radius * math.sin(angle))

        bx = px - Inches(1.0)
        by = py - Inches(0.55)
        bw = Inches(2.0)
        bh = Inches(1.1)

        phase_box = slide.shapes.add_shape(
            5, bx, by, bw, bh
        )
        phase_box.fill.solid()
        phase_box.fill.fore_color.rgb = theme.WHITE_RGB
        phase_box.line.color.rgb = theme.BORDER_RGB
        phase_box.line.width = Pt(0.5)

        sp = phase_box._element
        spPr = sp.find(qn('p:spPr'))
        if spPr is not None:
            effLst = spPr.find(qn('a:effectLst'))
            if effLst is None:
                effLst = etree.SubElement(spPr, qn('a:effectLst'))
            outerShdw = etree.SubElement(effLst, qn('a:outerShdw'))
            outerShdw.set('blurRad', str(Emu(40000)))
            outerShdw.set('dist', str(Emu(15240)))
            outerShdw.set('dir', str(18000000))
            outerShdw.set('algn', 'tl')
            srgbClr = etree.SubElement(outerShdw, qn('a:srgbClr'))
            srgbClr.set('val', '000000')
            alpha = etree.SubElement(srgbClr, qn('a:alpha'))
            alpha.set('val', '12000')

        tbox = slide.shapes.add_textbox(bx + Inches(0.15), by + Inches(0.1), bw - Inches(0.3), bh - Inches(0.2))
        ttf = tbox.text_frame
        ttf.word_wrap = True
        ttf.margin_left = Pt(0)
        ttf.margin_top = Pt(0)

        tp = ttf.paragraphs[0]
        tp.alignment = PP_ALIGN.CENTER
        tp.text = p.get("title", "")
        tp.font.size = Pt(14)
        tp.font.bold = True
        tp.font.color.rgb = theme.DARK_RGB
        tp.font.name = theme.FONT_FAMILY

        if p.get("text"):
            dp = ttf.add_paragraph()
            dp.alignment = PP_ALIGN.CENTER
            dp.text = p["text"]
            dp.font.size = Pt(12)
            dp.font.color.rgb = theme.MEDIUM_RGB
            dp.font.name = theme.FONT_FAMILY
            dp.space_before = Pt(2)

    center = slide.shapes.add_shape(
        9, cx - Inches(0.7), cy - Inches(0.7),
        Inches(1.4), Inches(1.4)
    )
    center.fill.solid()
    center.fill.fore_color.rgb = theme.BANNER_ORANGE
    center.line.fill.background()

    ctbox = slide.shapes.add_textbox(
        cx - Inches(0.6), cy - Inches(0.35),
        Inches(1.2), Inches(0.7)
    )
    ctf = ctbox.text_frame
    ctf.word_wrap = True
    ctf.margin_left = Pt(0)
    ctf.margin_top = Pt(0)
    cp = ctf.paragraphs[0]
    cp.alignment = PP_ALIGN.CENTER
    cp.text = "Utilisateur"
    cp.font.size = Pt(14)
    cp.font.bold = True
    cp.font.color.rgb = theme.WHITE_RGB
    cp.font.name = theme.FONT_FAMILY


def _add_split_premium(slide, left_text, right_text, left_title="", right_title=""):
    mid = Inches(0.08)
    half_w = (theme.CONTENT_W - mid) // 2
    y = Inches(2.2)
    h = Inches(3.6)

    left_box = slide.shapes.add_shape(
        5, theme.MARGIN_L, y, half_w, h
    )
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = theme.LIGHT_BG_RGB
    left_box.line.color.rgb = theme.BORDER_RGB
    left_box.line.width = Pt(0.5)

    if left_title:
        tbox = slide.shapes.add_textbox(
            theme.MARGIN_L + Inches(0.35), y + Inches(0.2),
            half_w - Inches(0.7), Inches(0.45)
        )
        ttf = tbox.text_frame
        ttf.margin_left = Pt(0)
        ttf.margin_top = Pt(0)
        tp = ttf.paragraphs[0]
        tp.text = left_title
        tp.font.size = Pt(18)
        tp.font.bold = True
        tp.font.color.rgb = theme.BANNER_ORANGE
        tp.font.name = theme.FONT_FAMILY

    bbox = slide.shapes.add_textbox(
        theme.MARGIN_L + Inches(0.35), y + Inches(0.7),
        half_w - Inches(0.7), h - Inches(0.9)
    )
    btf = bbox.text_frame
    btf.word_wrap = True
    btf.margin_left = Pt(0)
    btf.margin_top = Pt(0)
    bp = btf.paragraphs[0]
    bp.text = left_text
    bp.font.size = theme.BODY_SIZE
    bp.font.color.rgb = theme.DARK_RGB
    bp.font.name = theme.FONT_FAMILY

    arrow_box = slide.shapes.add_shape(
        1, theme.MARGIN_L + half_w + Inches(0.02),
        y + h // 2 - Inches(0.4), Inches(0.04), Inches(0.8)
    )
    arrow_box.fill.solid()
    arrow_box.fill.fore_color.rgb = theme.BANNER_ORANGE
    arrow_box.line.fill.background()

    right_x = theme.MARGIN_L + half_w + mid
    right_box = slide.shapes.add_shape(
        5, right_x, y, half_w, h
    )
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = theme.LIGHT_BG_RGB
    right_box.line.color.rgb = theme.BORDER_RGB
    right_box.line.width = Pt(0.5)

    if right_title:
        tbox2 = slide.shapes.add_textbox(
            right_x + Inches(0.35), y + Inches(0.2),
            half_w - Inches(0.7), Inches(0.45)
        )
        ttf2 = tbox2.text_frame
        ttf2.margin_left = Pt(0)
        ttf2.margin_top = Pt(0)
        tp2 = ttf2.paragraphs[0]
        tp2.text = right_title
        tp2.font.size = Pt(18)
        tp2.font.bold = True
        tp2.font.color.rgb = theme.BANNER_ORANGE
        tp2.font.name = theme.FONT_FAMILY

    bbox2 = slide.shapes.add_textbox(
        right_x + Inches(0.35), y + Inches(0.7),
        half_w - Inches(0.7), h - Inches(0.9)
    )
    btf2 = bbox2.text_frame
    btf2.word_wrap = True
    btf2.margin_left = Pt(0)
    btf2.margin_top = Pt(0)
    bp2 = btf2.paragraphs[0]
    bp2.text = right_text
    bp2.font.size = theme.BODY_SIZE
    bp2.font.color.rgb = theme.DARK_RGB
    bp2.font.name = theme.FONT_FAMILY


def build_cover(slide, deck):
    _add_bg(slide, theme.WHITE_RGB)

    band = slide.shapes.add_shape(
        1, Inches(0.3), Inches(0), Inches(0.15), theme.SLIDE_H
    )
    band.fill.solid()
    band.fill.fore_color.rgb = theme.BANNER_ORANGE
    band.line.fill.background()

    decor.add_top_accent_line(slide)
    decor.add_geometric_dot(slide)

    title = deck.get("deckTitle", "Presentation")
    box = slide.shapes.add_textbox(
        Inches(1.5), Inches(2.0), Inches(10.0), Inches(1.6)
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    tf.margin_top = Pt(0)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = theme.DARK_RGB
    p.font.name = theme.FONT_FAMILY

    if deck.get("deckSubtitle"):
        sbox = slide.shapes.add_textbox(
            Inches(1.5), Inches(3.6), Inches(10.0), Inches(0.6)
        )
        stf = sbox.text_frame
        stf.margin_left = Pt(0)
        stf.margin_top = Pt(0)
        sp = stf.paragraphs[0]
        sp.text = deck["deckSubtitle"]
        sp.font.size = Pt(22)
        sp.font.color.rgb = theme.MEDIUM_RGB
        sp.font.name = theme.FONT_FAMILY

    date_auth = []
    if deck.get("date"):
        date_auth.append(deck["date"])
    if deck.get("authors"):
        date_auth.append(deck["authors"])
    if date_auth:
        dbox = slide.shapes.add_textbox(
            Inches(1.5), Inches(4.5), Inches(10.0), Inches(0.5)
        )
        dtf = dbox.text_frame
        dtf.margin_left = Pt(0)
        dtf.margin_top = Pt(0)
        dp = dtf.paragraphs[0]
        dp.text = " | ".join(date_auth)
        dp.font.size = Pt(14)
        dp.font.color.rgb = theme.MEDIUM_RGB
        dp.font.name = theme.FONT_FAMILY

    footer.add_footer(slide)

    return {"cover": True}


def build_message_only(slide, slide_data):
    _add_bg(slide, theme.WHITE_RGB)
    decor.add_design_elements(slide, 2)

    if slide_data.get("title"):
        _add_title(slide, slide_data["title"])

    main_msg = slide_data.get("mainMessage", slide_data.get("message", ""))
    if main_msg:
        box = slide.shapes.add_textbox(
            theme.MARGIN_L, Inches(2.5), theme.CONTENT_W, Inches(2.5)
        )
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(0)
        tf.margin_top = Pt(0)
        p = tf.paragraphs[0]
        p.text = main_msg
        p.font.size = Pt(28)
        p.font.color.rgb = theme.DARK_RGB
        p.font.name = theme.FONT_FAMILY

    footer.add_footer(slide)
    if slide_data.get("conclusion"):
        conc.add_conclusion_banner(slide, slide_data["conclusion"])


def build_three_cards(slide, slide_data):
    _add_bg(slide, theme.WHITE_RGB)
    decor.add_design_elements(slide, 2)

    if slide_data.get("title"):
        _add_title(slide, slide_data["title"])
    if slide_data.get("subtitle"):
        _add_subtitle(slide, slide_data["subtitle"])

    cards = slide_data.get("cards", [])
    _add_card_grid(slide, cards[:3], cols=3)

    footer.add_footer(slide)
    if slide_data.get("conclusion"):
        conc.add_conclusion_banner(slide, slide_data["conclusion"])


def build_four_cards(slide, slide_data):
    _add_bg(slide, theme.WHITE_RGB)
    decor.add_design_elements(slide, 2)

    if slide_data.get("title"):
        _add_title(slide, slide_data["title"])
    if slide_data.get("subtitle"):
        _add_subtitle(slide, slide_data["subtitle"])

    cards = slide_data.get("cards", [])
    _add_card_grid(slide, cards[:4], cols=2)

    footer.add_footer(slide)
    if slide_data.get("conclusion"):
        conc.add_conclusion_banner(slide, slide_data["conclusion"])


def build_problem_solution(slide, slide_data):
    _add_bg(slide, theme.WHITE_RGB)
    decor.add_design_elements(slide, 2)

    if slide_data.get("title"):
        _add_title(slide, slide_data["title"])

    _add_split_premium(
        slide,
        slide_data.get("problem", ""),
        slide_data.get("solution", ""),
        left_title="Problème",
        right_title="Solution"
    )

    footer.add_footer(slide)
    if slide_data.get("conclusion"):
        conc.add_conclusion_banner(slide, slide_data["conclusion"])


def build_process(slide, slide_data):
    _add_bg(slide, theme.WHITE_RGB)
    decor.add_design_elements(slide, 2)

    if slide_data.get("title"):
        _add_title(slide, slide_data["title"])

    steps = slide_data.get("steps", slide_data.get("cards", []))
    _add_step_line_premium(slide, steps[:6])

    footer.add_footer(slide)
    if slide_data.get("conclusion"):
        conc.add_conclusion_banner(slide, slide_data["conclusion"])


def build_kpi(slide, slide_data):
    _add_bg(slide, theme.WHITE_RGB)
    decor.add_design_elements(slide, 2)

    if slide_data.get("title"):
        _add_title(slide, slide_data["title"])
    if slide_data.get("subtitle"):
        _add_subtitle(slide, slide_data["subtitle"])

    kpis = slide_data.get("kpis", [])
    n = min(len(kpis), 4)
    if n > 0:
        kpi_w = theme.CONTENT_W // n
        y = Inches(2.6)
        for i, k in enumerate(kpis):
            _add_kpi_big(
                slide, k.get("value", ""), k.get("label", ""),
                theme.MARGIN_L + kpi_w * i, y, kpi_w
            )

    footer.add_footer(slide)
    if slide_data.get("conclusion"):
        conc.add_conclusion_banner(slide, slide_data["conclusion"])


def build_lessons(slide, slide_data):
    _add_bg(slide, theme.WHITE_RGB)
    decor.add_design_elements(slide, 2)

    if slide_data.get("title"):
        _add_title(slide, slide_data["title"])

    lessons = slide_data.get("lessons", slide_data.get("cards", []))
    y = Inches(2.3)
    lesson_w = Inches(10.5)
    lesson_h = Inches(0.7)
    gap = Inches(0.18)

    for i, l in enumerate(lessons[:5]):
        top = y + i * (lesson_h + gap)
        text = l if isinstance(l, str) else l.get("text", l.get("title", ""))

        num_bg = slide.shapes.add_shape(
            9, theme.MARGIN_L, top + Inches(0.12),
            Inches(0.4), Inches(0.4)
        )
        num_bg.fill.solid()
        num_bg.fill.fore_color.rgb = theme.BANNER_ORANGE
        num_bg.line.fill.background()

        nbox = slide.shapes.add_textbox(
            theme.MARGIN_L, top + Inches(0.12),
            Inches(0.4), Inches(0.4)
        )
        ntf = nbox.text_frame
        np = ntf.paragraphs[0]
        np.alignment = PP_ALIGN.CENTER
        np.text = str(i + 1)
        np.font.size = Pt(14)
        np.font.bold = True
        np.font.color.rgb = theme.WHITE_RGB
        np.font.name = theme.FONT_FAMILY
        ntf.margin_left = Pt(0)
        ntf.margin_top = Pt(2)

        tbox = slide.shapes.add_textbox(
            theme.MARGIN_L + Inches(0.6), top,
            lesson_w - Inches(0.6), lesson_h
        )
        ttf = tbox.text_frame
        ttf.word_wrap = True
        ttf.margin_left = Pt(0)
        ttf.margin_top = Pt(0)
        tp = ttf.paragraphs[0]
        tp.text = text
        tp.font.size = Pt(16)
        tp.font.color.rgb = theme.DARK_RGB
        tp.font.name = theme.FONT_FAMILY

        if i < len(lessons) - 1:
            sep = slide.shapes.add_shape(
                1, theme.MARGIN_L + Inches(0.6),
                int(top + lesson_h), lesson_w - Inches(0.6), Pt(0.5)
            )
            sep.fill.solid()
            sep.fill.fore_color.rgb = theme.BORDER_RGB
            sep.line.fill.background()

    footer.add_footer(slide)
    if slide_data.get("conclusion"):
        conc.add_conclusion_banner(slide, slide_data["conclusion"])


def build_role_focus(slide, slide_data):
    _add_bg(slide, theme.WHITE_RGB)
    decor.add_design_elements(slide, 2)

    if slide_data.get("title"):
        _add_title(slide, slide_data["title"])

    role_name = slide_data.get("roleName", "Role")
    role_title = slide_data.get("roleTitle", "")
    responsibilities = slide_data.get("responsibilities", slide_data.get("cards", []))

    cx = int(Inches(6.666))
    cy = int(Inches(3.0))

    center = slide.shapes.add_shape(
        9, cx - Inches(0.6), cy - Inches(0.6),
        Inches(1.2), Inches(1.2)
    )
    center.fill.solid()
    center.fill.fore_color.rgb = theme.BANNER_ORANGE
    center.line.fill.background()

    cbox = slide.shapes.add_textbox(
        cx - Inches(0.55), cy - Inches(0.4),
        Inches(1.1), Inches(0.8)
    )
    ctf = cbox.text_frame
    ctf.word_wrap = True
    ctf.margin_left = Pt(0)
    ctf.margin_top = Pt(0)
    cp = ctf.paragraphs[0]
    cp.alignment = PP_ALIGN.CENTER
    cp.text = role_name
    cp.font.size = Pt(14)
    cp.font.bold = True
    cp.font.color.rgb = theme.WHITE_RGB
    cp.font.name = theme.FONT_FAMILY

    if role_title:
        sp = ctf.add_paragraph()
        sp.alignment = PP_ALIGN.CENTER
        sp.text = role_title
        sp.font.size = Pt(10)
        sp.font.color.rgb = theme.LIGHT_BG_RGB
        sp.font.name = theme.FONT_FAMILY
        sp.space_before = Pt(2)

    n = min(len(responsibilities), 4)
    import math
    for i, r in enumerate(responsibilities[:n]):
        angle = 2 * math.pi * i / n - math.pi / 2
        px = int(cx + Inches(2.0) * math.cos(angle))
        py = int(cy + Inches(1.5) * math.sin(angle))

        text = r if isinstance(r, str) else r.get("text", r.get("title", ""))
        rbox = slide.shapes.add_textbox(
            px - Inches(0.8), py - Inches(0.2),
            Inches(1.6), Inches(0.6)
        )
        rtf = rbox.text_frame
        rtf.word_wrap = True
        rtf.margin_left = Pt(0)
        rtf.margin_top = Pt(0)
        rp = rtf.paragraphs[0]
        rp.alignment = PP_ALIGN.CENTER
        rp.text = text
        rp.font.size = Pt(14)
        rp.font.color.rgb = theme.DARK_RGB
        rp.font.name = theme.FONT_FAMILY

    footer.add_footer(slide)
    if slide_data.get("conclusion"):
        conc.add_conclusion_banner(slide, slide_data["conclusion"])


def build_adoption_loop(slide, slide_data):
    _add_bg(slide, theme.WHITE_RGB)
    decor.add_design_elements(slide, 2)

    if slide_data.get("title"):
        _add_title(slide, slide_data["title"])

    phases = slide_data.get("phases", slide_data.get("cards", []))
    _add_cycle_premium(slide, phases[:4])

    footer.add_footer(slide)
    if slide_data.get("conclusion"):
        conc.add_conclusion_banner(slide, slide_data["conclusion"])


def build_closing(slide, slide_data):
    _add_bg(slide, theme.WHITE_RGB)
    decor.add_design_elements(slide, 2)

    main_msg = slide_data.get("mainMessage", slide_data.get("message", ""))
    if main_msg:
        box = slide.shapes.add_textbox(
            theme.MARGIN_L, Inches(2.0), theme.CONTENT_W, Inches(1.5)
        )
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(0)
        tf.margin_top = Pt(0)
        p = tf.paragraphs[0]
        p.text = main_msg
        p.font.size = Pt(34)
        p.font.bold = True
        p.font.color.rgb = theme.DARK_RGB
        p.font.name = theme.FONT_FAMILY

    takeaways = slide_data.get("takeaways", slide_data.get("cards", []))
    y = Inches(3.8)
    for i, t in enumerate(takeaways[:3]):
        text = t if isinstance(t, str) else t.get("text", t.get("title", ""))
        tbox = slide.shapes.add_textbox(
            theme.MARGIN_L + Inches(0.3), y + i * Inches(0.7),
            theme.CONTENT_W - Inches(0.3), Inches(0.55)
        )
        ttf = tbox.text_frame
        ttf.word_wrap = True
        ttf.margin_left = Pt(0)
        ttf.margin_top = Pt(0)
        tp = ttf.paragraphs[0]
        run = tp.add_run()
        run.text = f"→  {text}"
        run.font.size = Pt(18)
        run.font.color.rgb = theme.DARK_RGB
        run.font.name = theme.FONT_FAMILY

    footer.add_footer(slide)
    if slide_data.get("conclusion"):
        conc.add_conclusion_banner(slide, slide_data["conclusion"])


BUILDERS = {
    "cover_orange": build_cover,
    "message_only": build_message_only,
    "three_cards": build_three_cards,
    "four_cards": build_four_cards,
    "problem_solution": build_problem_solution,
    "process_horizontal": build_process,
    "kpi_context": build_kpi,
    "lessons_learned": build_lessons,
    "role_focus": build_role_focus,
    "adoption_loop": build_adoption_loop,
    "closing": build_closing,
}
